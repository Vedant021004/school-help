import os
from pathlib import Path
from typing import Optional

from backend.config import settings
from backend.models import QuestionPaper, AnswerKey


# ==========================================
# 1. REPORTLAB PDF EXPORTERS
# ==========================================

def export_question_paper_pdf(paper: QuestionPaper) -> str:
    """Generates a professional examination PDF for the Question Paper."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, KeepTogether
    )

    filename = f"Question_Paper_{paper.id[:8]}.pdf"
    file_path = settings.EXPORTS_DIR / filename

    doc = SimpleDocTemplate(
        str(file_path),
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'ExamTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=14,
        alignment=1, # Center
        spaceAfter=4
    )
    subtitle_style = ParagraphStyle(
        'ExamSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=11,
        alignment=1,
        spaceAfter=6
    )
    meta_style = ParagraphStyle(
        'ExamMeta',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        leading=12
    )
    sec_heading_style = ParagraphStyle(
        'SectionHeading',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=11,
        textColor=colors.HexColor("#1e3a8a"),
        spaceBefore=10,
        spaceAfter=4
    )
    q_style = ParagraphStyle(
        'QuestionText',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=13,
        spaceBefore=4,
        spaceAfter=3
    )
    opt_style = ParagraphStyle(
        'OptionText',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        leading=12,
        leftIndent=15
    )

    elements = []

    # School Header
    elements.append(Paragraph(paper.school_name.upper(), title_style))
    elements.append(Paragraph(f"{paper.exam_name} • Academic Year {paper.date_str or '2025-26'}", subtitle_style))
    elements.append(Spacer(1, 4))

    # Meta Table: Class, Subject, Max Marks, Time
    meta_data = [
        [
            Paragraph(f"<b>Subject:</b> {paper.subject}", meta_style),
            Paragraph(f"<b>Class:</b> {paper.grade} ({paper.board})", meta_style),
            Paragraph(f"<b>Max Marks:</b> {paper.total_marks}", meta_style),
            Paragraph(f"<b>Time Allowed:</b> {paper.duration_minutes // 60} Hours", meta_style)
        ]
    ]
    meta_table = Table(meta_data, colWidths=[140, 130, 110, 130])
    meta_table.setStyle(TableStyle([
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor("#94a3b8")),
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor("#f8fafc")),
        ('TOPPADDING', (0,0), (-1,-1), 5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 5),
    ]))
    elements.append(meta_table)
    elements.append(Spacer(1, 8))

    # General Instructions
    if paper.instructions:
        elements.append(Paragraph("<b>General Instructions:</b>", ParagraphStyle('InstHead', fontName='Helvetica-Bold', fontSize=9, spaceAfter=2)))
        for idx, inst in enumerate(paper.instructions):
            elements.append(Paragraph(f"{idx+1}. {inst}", ParagraphStyle('InstItem', fontName='Helvetica', fontSize=8.5, leading=11)))
        elements.append(Spacer(1, 6))

    elements.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#cbd5e1"), spaceAfter=10))

    # Group questions by section
    sections_map = {}
    for q in paper.questions:
        if q.section_name not in sections_map:
            sections_map[q.section_name] = []
        sections_map[q.section_name].append(q)

    for sec_name, q_list in sections_map.items():
        # Section Header
        sec_marks = sum(q.marks for q in q_list)
        elements.append(Paragraph(f"<b>{sec_name}</b> ({len(q_list)} Questions • {sec_marks} Marks)", sec_heading_style))

        for q in q_list:
            q_flowables = []
            # Question and Marks line
            q_row = [
                [
                    Paragraph(f"<b>Q{q.question_number}.</b> {q.question_text}", q_style),
                    Paragraph(f"<b>[{q.marks}]</b>", ParagraphStyle('QMarks', fontName='Helvetica-Bold', fontSize=9.5, alignment=2))
                ]
            ]
            q_table = Table(q_row, colWidths=[460, 50])
            q_table.setStyle(TableStyle([
                ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ('LEFTPADDING', (0,0), (-1,-1), 0),
                ('RIGHTPADDING', (0,0), (-1,-1), 0),
            ]))
            q_flowables.append(q_table)

            # Options for MCQs
            if q.options:
                for opt in q.options:
                    q_flowables.append(Paragraph(opt, opt_style))
                q_flowables.append(Spacer(1, 3))

            q_flowables.append(Spacer(1, 4))
            elements.append(KeepTogether(q_flowables))

    elements.append(Spacer(1, 15))
    elements.append(Paragraph("★★★ ALL THE BEST ★★★", ParagraphStyle('FooterEnd', fontName='Helvetica-Bold', fontSize=10, alignment=1)))

    doc.build(elements)
    return str(file_path)


def export_answer_key_pdf(ak: AnswerKey) -> str:
    """Generates a professional examination PDF for the Answer Key."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, KeepTogether
    )

    filename = f"Answer_Key_{ak.paper_id[:8]}.pdf"
    file_path = settings.EXPORTS_DIR / filename

    doc = SimpleDocTemplate(
        str(file_path),
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('AKTitle', parent=styles['Heading1'], fontName='Helvetica-Bold', fontSize=14, alignment=1, spaceAfter=4)
    sub_style = ParagraphStyle('AKSub', parent=styles['Normal'], fontName='Helvetica-Bold', fontSize=10, alignment=1, spaceAfter=8)
    q_style = ParagraphStyle('AKQ', fontName='Helvetica-Bold', fontSize=9.5, textColor=colors.HexColor("#0f172a"), spaceBefore=6)
    ans_style = ParagraphStyle('AKAns', fontName='Helvetica', fontSize=9, leading=12, textColor=colors.HexColor("#047857"))
    source_style = ParagraphStyle('AKSource', fontName='Helvetica-Oblique', fontSize=8, textColor=colors.HexColor("#64748b"), spaceAfter=6)

    elements = [
        Paragraph("OFFICIAL MARKING SCHEME & ANSWER KEY", title_style),
        Paragraph(f"Paper: {ak.paper_title} • Subject: {ak.subject} ({ak.grade}) • Total: {ak.total_marks} Marks", sub_style),
        HRFlowable(width="100%", thickness=1, color=colors.HexColor("#059669"), spaceAfter=10)
    ]

    for item in ak.answers:
        item_flows = []
        item_flows.append(Paragraph(f"<b>Q{item.question_number} [{item.marks} Mark{'s' if item.marks > 1 else ''}]</b> - <i>{item.question_text[:100]}...</i>", q_style))
        item_flows.append(Paragraph(f"<b>Correct Answer / Key:</b> {item.correct_answer}", ans_style))

        if item.formula_and_steps:
            item_flows.append(Paragraph(f"<b>Formula & Steps:</b> {item.formula_and_steps}", ParagraphStyle('Steps', fontName='Helvetica', fontSize=8.5, leading=11)))

        if item.detailed_explanation:
            item_flows.append(Paragraph(f"<b>Explanation:</b> {item.detailed_explanation}", ParagraphStyle('Exp', fontName='Helvetica', fontSize=8.5, leading=11)))

        # Citation
        src = item.source_reference
        item_flows.append(Paragraph(f"📚 <b>Textbook Grounding:</b> {src.book_title} | Chapter: {src.chapter_name} (Page {src.page})", source_style))
        item_flows.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#e2e8f0"), spaceAfter=4))
        elements.append(KeepTogether(item_flows))

    doc.build(elements)
    return str(file_path)


# ==========================================
# 2. PYTHON-DOCX EXPORTERS
# ==========================================

def export_question_paper_docx(paper: QuestionPaper) -> str:
    """Generates a styled Microsoft Word (.docx) Question Paper."""
    import docx
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    filename = f"Question_Paper_{paper.id[:8]}.docx"
    file_path = settings.EXPORTS_DIR / filename

    document = docx.Document()

    # School Title Header
    p_title = document.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_title = p_title.add_run(paper.school_name.upper())
    run_title.font.name = 'Calibri'
    run_title.font.size = Pt(16)
    run_title.font.bold = True

    p_sub = document.add_paragraph()
    p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_sub = p_sub.add_run(f"{paper.exam_name} • {paper.date_str or '2025-26'}\nSubject: {paper.subject} | Grade: {paper.grade} | Time: {paper.duration_minutes // 60} Hours | Max Marks: {paper.total_marks}")
    run_sub.font.size = Pt(11)

    # General Instructions
    if paper.instructions:
        document.add_heading("General Instructions:", level=2)
        for idx, inst in enumerate(paper.instructions):
            document.add_paragraph(f"{idx+1}. {inst}", style='List Number' if 'List Number' in document.styles else None)

    # Sections & Questions
    sections_map = {}
    for q in paper.questions:
        if q.section_name not in sections_map:
            sections_map[q.section_name] = []
        sections_map[q.section_name].append(q)

    for sec_name, q_list in sections_map.items():
        document.add_heading(f"{sec_name} ({sum(q.marks for q in q_list)} Marks)", level=1)
        for q in q_list:
            p_q = document.add_paragraph()
            r_qnum = p_q.add_run(f"Q{q.question_number}. ")
            r_qnum.bold = True
            p_q.add_run(f"{q.question_text} ")
            r_marks = p_q.add_run(f"[{q.marks} Mark{'s' if q.marks > 1 else ''}]")
            r_marks.bold = True

            if q.options:
                for opt in q.options:
                    p_opt = document.add_paragraph(opt)
                    p_opt.paragraph_format.left_indent = Inches(0.3)

    document.save(str(file_path))
    return str(file_path)


# ==========================================
# 3. PYTHON-PPTX PRESENTATION EXPORTER
# ==========================================

def export_presentation_pptx(deck: 'SlideDeck') -> str:
    """Generates a professional, styled Microsoft PowerPoint (.pptx) presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    filename = f"Presentation_{deck.id[:8]}.pptx"
    file_path = settings.EXPORTS_DIR / filename

    prs = Presentation()
    # Set slide width and height to 16:9 widescreen (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    PRIMARY_COLOR = RGBColor(30, 58, 138)     # Deep Indigo
    ACCENT_COLOR = RGBColor(79, 70, 229)      # Bright Indigo
    TEXT_COLOR = RGBColor(30, 41, 59)         # Dark Slate
    MUTED_COLOR = RGBColor(100, 116, 139)     # Slate Grey
    BG_LIGHT = RGBColor(248, 250, 252)        # Light Slate

    blank_slide_layout = prs.slide_layouts[6]

    for slide_data in deck.slides:
        slide = prs.slides.add_slide(blank_slide_layout)

        # Title Slide
        if slide_data.layout == "title":
            # Main Title Box
            tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(2.5))
            tf = tx_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.name = 'Calibri'
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
            p.alignment = PP_ALIGN.CENTER

            # Subtitle Box
            p_sub = tf.add_paragraph()
            p_sub.text = deck.subtitle or f"{deck.grade} • {deck.subject} • NCERT Master Class"
            p_sub.font.size = Pt(22)
            p_sub.font.color.rgb = ACCENT_COLOR
            p_sub.alignment = PP_ALIGN.CENTER
            p_sub.space_before = Pt(14)

            # Meta Bullet points
            if slide_data.bullet_points:
                p_meta = tf.add_paragraph()
                p_meta.text = " • ".join(slide_data.bullet_points)
                p_meta.font.size = Pt(14)
                p_meta.font.color.rgb = MUTED_COLOR
                p_meta.alignment = PP_ALIGN.CENTER
                p_meta.space_before = Pt(18)

        else:
            # Content Slide Header Banner
            header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
            tf_h = header_box.text_frame
            tf_h.word_wrap = True
            p_h = tf_h.paragraphs[0]
            p_h.text = slide_data.title
            p_h.font.name = 'Calibri'
            p_h.font.size = Pt(30)
            p_h.font.bold = True
            p_h.font.color.rgb = PRIMARY_COLOR

            p_h_sub = tf_h.add_paragraph()
            p_h_sub.text = f"{deck.subject} • {deck.chapter_name}"
            p_h_sub.font.size = Pt(12)
            p_h_sub.font.color.rgb = MUTED_COLOR

            # Content Left/Main Box
            width = Inches(7.5) if (slide_data.key_definition or slide_data.activity_box) else Inches(11.5)
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), width, Inches(5.0))
            tf_c = content_box.text_frame
            tf_c.word_wrap = True

            for idx, bp in enumerate(slide_data.bullet_points):
                p_bp = tf_c.paragraphs[0] if idx == 0 else tf_c.add_paragraph()
                p_bp.text = f"•  {bp}"
                p_bp.font.name = 'Calibri'
                p_bp.font.size = Pt(18)
                p_bp.font.color.rgb = TEXT_COLOR
                p_bp.space_after = Pt(12)

            # Side Box (Key Definition or Activity)
            if slide_data.key_definition or slide_data.activity_box:
                side_box = slide.shapes.add_textbox(Inches(8.7), Inches(1.8), Inches(3.8), Inches(4.8))
                tf_s = side_box.text_frame
                tf_s.word_wrap = True

                p_st = tf_s.paragraphs[0]
                p_st.text = "📌 KEY TAKEAWAY" if slide_data.key_definition else "⚡ CLASS ACTIVITY"
                p_st.font.name = 'Calibri'
                p_st.font.size = Pt(15)
                p_st.font.bold = True
                p_st.font.color.rgb = ACCENT_COLOR
                p_st.space_after = Pt(8)

                p_sc = tf_s.add_paragraph()
                p_sc.text = slide_data.key_definition or slide_data.activity_box or ""
                p_sc.font.name = 'Calibri'
                p_sc.font.size = Pt(14)
                p_sc.font.color.rgb = TEXT_COLOR

        # Speaker notes
        if slide_data.speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data.speaker_notes

    prs.save(str(file_path))
    return str(file_path)


# ==========================================
# 4. WORKSHEET PDF EXPORTER
# ==========================================

def export_worksheet_pdf(worksheet: 'ChapterWorksheet', school_name: str = "Central Academy") -> str:
    """Generates a clean, printable classroom PDF worksheet with answer spaces."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, KeepTogether
    )

    filename = f"Worksheet_{worksheet.id[:8]}.pdf"
    file_path = settings.EXPORTS_DIR / filename

    doc = SimpleDocTemplate(
        str(file_path),
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'WsTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=13,
        alignment=1,
        spaceAfter=3
    )
    subtitle_style = ParagraphStyle(
        'WsSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=10,
        alignment=1,
        spaceAfter=6
    )
    meta_style = ParagraphStyle(
        'WsMeta',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        leading=12
    )
    q_style = ParagraphStyle(
        'WsQuestion',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=13,
        spaceBefore=4,
        spaceAfter=3
    )
    opt_style = ParagraphStyle(
        'WsOption',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        leading=12,
        leftIndent=15
    )

    elements = []

    # School & Worksheet Title
    elements.append(Paragraph(school_name.upper(), title_style))
    elements.append(Paragraph(worksheet.title, subtitle_style))
    elements.append(Spacer(1, 4))

    # Student Info Table
    info_data = [
        [
            Paragraph("<b>Student Name:</b> ___________________________", meta_style),
            Paragraph("<b>Roll No:</b> ____________", meta_style),
            Paragraph("<b>Date:</b> ____________", meta_style)
        ],
        [
            Paragraph(f"<b>Class & Section:</b> {worksheet.grade}", meta_style),
            Paragraph(f"<b>Subject:</b> {worksheet.subject}", meta_style),
            Paragraph(f"<b>Max Marks:</b> {worksheet.total_marks} | <b>Time:</b> {worksheet.estimated_time_minutes}m", meta_style)
        ]
    ]
    info_table = Table(info_data, colWidths=[200, 160, 160])
    info_table.setStyle(TableStyle([
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor("#94a3b8")),
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor("#f8fafc")),
        ('TOPPADDING', (0,0), (-1,-1), 4),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
    ]))
    elements.append(info_table)
    elements.append(Spacer(1, 8))

    # Instructions
    if worksheet.instructions:
        elements.append(Paragraph("<b>Instructions:</b> " + " • ".join(worksheet.instructions), ParagraphStyle('WsInst', fontName='Helvetica', fontSize=8, leading=10, textColor=colors.HexColor("#475569"))))
        elements.append(Spacer(1, 4))

    elements.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#cbd5e1"), spaceAfter=8))

    # Questions
    for q in worksheet.questions:
        q_flows = []
        q_row = [
            [
                Paragraph(f"<b>Q{q.question_number}.</b> {q.question_text}", q_style),
                Paragraph(f"<b>[{q.marks}M]</b>", ParagraphStyle('QM', fontName='Helvetica-Bold', fontSize=9, alignment=2))
            ]
        ]
        q_table = Table(q_row, colWidths=[465, 45])
        q_table.setStyle(TableStyle([
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('LEFTPADDING', (0,0), (-1,-1), 0),
            ('RIGHTPADDING', (0,0), (-1,-1), 0),
        ]))
        q_flows.append(q_table)

        if q.options:
            for opt in q.options:
                q_flows.append(Paragraph(opt, opt_style))
            q_flows.append(Spacer(1, 4))
        else:
            # Leave blank answer lines based on marks
            line_count = 2 if q.marks <= 2 else 4 if q.marks <= 3 else 7
            for _ in range(line_count):
                q_flows.append(HRFlowable(width="100%", thickness=0.3, color=colors.HexColor("#e2e8f0"), spaceBefore=8, spaceAfter=4))

        q_flows.append(Spacer(1, 6))
        elements.append(KeepTogether(q_flows))

    doc.build(elements)
    return str(file_path)


# ==========================================
# 5. CHAPTER NOTES PDF EXPORTER
# ==========================================

def export_chapter_notes_pdf(notes: 'ChapterNotes') -> str:
    """Generates a styled revision notes PDF document."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, KeepTogether
    )

    filename = f"Notes_{notes.chapter_id[:8]}.pdf"
    file_path = settings.EXPORTS_DIR / filename

    doc = SimpleDocTemplate(
        str(file_path),
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('NTitle', parent=styles['Heading1'], fontName='Helvetica-Bold', fontSize=14, alignment=1, spaceAfter=4)
    sub_style = ParagraphStyle('NSub', parent=styles['Normal'], fontName='Helvetica-Bold', fontSize=10, alignment=1, textColor=colors.HexColor("#4338ca"), spaceAfter=8)
    h2_style = ParagraphStyle('NH2', parent=styles['Heading2'], fontName='Helvetica-Bold', fontSize=11, textColor=colors.HexColor("#1e3a8a"), spaceBefore=8, spaceAfter=4)
    body_style = ParagraphStyle('NBody', parent=styles['Normal'], fontName='Helvetica', fontSize=9, leading=13, spaceAfter=4)
    bullet_style = ParagraphStyle('NBullet', parent=styles['Normal'], fontName='Helvetica', fontSize=8.5, leading=12, leftIndent=12, spaceAfter=3)

    elements = []
    elements.append(Paragraph(notes.title, title_style))
    elements.append(Paragraph(f"Comprehensive NCERT High-Yield Study Notes • {notes.chapter_name}", sub_style))
    elements.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#cbd5e1"), spaceAfter=8))

    # Summary
    if notes.summary:
        elements.append(Paragraph("<b>Chapter Overview & Key Themes</b>", h2_style))
        elements.append(Paragraph(notes.summary, body_style))
        elements.append(Spacer(1, 6))

    # Definitions
    if notes.definitions:
        elements.append(Paragraph("<b>Key Terminology & Definitions</b>", h2_style))
        for d in notes.definitions:
            elements.append(Paragraph(f"• <b>{d.get('term', '')}:</b> {d.get('definition', '')}", bullet_style))
        elements.append(Spacer(1, 6))

    # Core Principles
    if notes.core_principles:
        elements.append(Paragraph("<b>Core Laws, Principles & Mechanisms</b>", h2_style))
        for cp in notes.core_principles:
            elements.append(Paragraph(f"• <b>{cp.get('title', '')}:</b> {cp.get('explanation', '')}", bullet_style))
        elements.append(Spacer(1, 6))

    # Formulas
    if notes.formulas:
        elements.append(Paragraph("<b>Governing Formulas & Equations</b>", h2_style))
        for f in notes.formulas:
            elements.append(Paragraph(f"• <b>{f.get('name', '')}:</b> <font color='#1e3a8a'><b>{f.get('formula', '')}</b></font> (Units: {f.get('units', 'SI Units')})", bullet_style))
        elements.append(Spacer(1, 6))

    # Revision Points
    if notes.revision_points:
        elements.append(Paragraph("<b>Final Revision & Exam Tips</b>", h2_style))
        for rp in notes.revision_points:
            elements.append(Paragraph(f"✓ {rp}", bullet_style))

    doc.build(elements)
    return str(file_path)


def export_textbook_chapter_pdf(book, chunks, chapter_id: Optional[str] = None) -> str:
    """
    Generates a full, real textbook PDF preview for any book/chapter.
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Table, TableStyle
    )

    filename = f"Textbook_{book.id}_{chapter_id or 'all'}.pdf"
    file_path = settings.EXPORTS_DIR / filename

    doc = SimpleDocTemplate(
        str(file_path),
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    styles = getSampleStyleSheet()
    header_style = ParagraphStyle(
        'TBHeader',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=18,
        leading=22,
        textColor=colors.HexColor('#1e1b4b'),
        spaceAfter=4
    )
    meta_style = ParagraphStyle(
        'TBMeta',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=10,
        textColor=colors.HexColor('#4f46e5'),
        spaceAfter=8
    )
    section_style = ParagraphStyle(
        'TBSec',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=12,
        leading=16,
        textColor=colors.HexColor('#0f172a'),
        spaceBefore=10,
        spaceAfter=4
    )
    body_style = ParagraphStyle(
        'TBBody',
        parent=styles['Normal'],
        fontName='Times-Roman',
        fontSize=10.5,
        leading=15,
        textColor=colors.HexColor('#1e293b'),
        spaceAfter=8
    )
    page_badge_style = ParagraphStyle(
        'TBPage',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=8.5,
        textColor=colors.HexColor('#64748b'),
        spaceAfter=2
    )

    elements = []
    elements.append(Paragraph(book.title, header_style))
    elements.append(Paragraph(f"NCERT Curriculum Edition • {book.grade or 'Standard'} • {book.subject or 'General'}", meta_style))
    elements.append(HRFlowable(width="100%", thickness=1.5, color=colors.HexColor("#4f46e5"), spaceAfter=12))

    if chunks:
        current_page = None
        for c in chunks:
            p_num = c.metadata.page_number if hasattr(c, 'metadata') else c.get('page_number', 1)
            sec_name = c.metadata.section_name if hasattr(c, 'metadata') else c.get('section_name', '')
            content = c.content if hasattr(c, 'content') else c.get('content', '')

            if p_num != current_page:
                current_page = p_num
                elements.append(Paragraph(f"<b>PAGE {p_num}</b>", page_badge_style))
                elements.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#cbd5e1"), spaceAfter=6))

            if sec_name and sec_name != "Main Excerpt":
                elements.append(Paragraph(sec_name, section_style))

            elements.append(Paragraph(content.replace("\n", "<br/>"), body_style))
            elements.append(Spacer(1, 4))
    else:
        elements.append(Paragraph("Textbook content excerpts indexed and verified by Shiksha-AI.", body_style))

    doc.build(elements)
    return str(file_path)


